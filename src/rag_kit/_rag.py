"""Main RAGSystem class — public API for rag-kit."""

from __future__ import annotations

import os
from typing import Any

from rag_kit import _docling
from rag_kit._llm import LLMConfig
from rag_kit._pipeline import Pipeline
from rag_kit._processor import (
    _build_section_mappings,
    _extract_headings_from_text,
    format_toc,
    process_chunks,
)
from rag_kit._search import search
from rag_kit._storage import Storage, compute_content_hash
from rag_kit._vector_index import VectorIndex

DEFAULT_CHUNK_SIZE = 1200
DEFAULT_CHUNK_OVERLAP = 200
DEFAULT_SEARCH_THRESHOLD = 0.6


def _has_meaningful_text(text: str, min_words: int = 10) -> bool:
    """Check if extracted text has enough real content for RAG.

    Scanned/image PDFs often return empty or gibberish text.
    """
    import re

    text = text.strip()
    if not text:
        return False
    words = re.findall(r"[A-Za-z]{2,}", text)
    if len(words) < min_words:
        return False
    tokens = text.split()
    if not tokens:
        return False
    return len(words) / len(tokens) > 0.3


def _clean_text(text: str) -> str:
    """Remove surrogate characters that crash UTF-8 encode/decode.

    Some PDF extractors produce surrogate characters (U+D800-U+DFFF)
    which Python's utf-8 codec refuses to encode.  This replaces them
    with '?' before any hash computation or chunk storage.
    """
    return text.encode("utf-8", errors="replace").decode("utf-8")


def _qnorm(text: str) -> str:
    """Normalize a question for cache keying: lowercase, collapse spaces,
    strip leading/trailing punctuation."""
    import re

    return re.sub(r"\s+", " ", text.strip().lower()).strip(" ?!.,;:")


class QueryResult:
    """Result of a query() call — answer text + citations + optional metrics."""

    def __init__(
        self, answer: str, citations: list[dict] | None = None, metrics: dict | None = None
    ):
        self.answer = answer
        self.citations = citations or []
        self.metrics = metrics or {}

    def __str__(self) -> str:
        return self.answer

    def __repr__(self) -> str:
        return f"QueryResult(answer={self.answer[:50]}..., citations={len(self.citations)})"


def _pdf_outline_headings(reader, page_texts: list[str]) -> list[dict]:
    """Convert a PDF's embedded outline (bookmarks) to heading dicts.

    Datasheets and manuals ship with real bookmarks; the chunk first-line
    heuristic produces register-name noise for them (e.g. Microchip
    datasheets → 2000+ junk headings). Outline items carry proper titles
    and hierarchy — page numbers become character offsets into the
    flattened page text so section mappings still line up with chunks.

    Returns [{title, level, offset}] in document order (already sorted).
    """
    outline = getattr(reader, "outline", None)
    if not outline:
        return []

    # Cumulative char offset of each page inside the "\n".join(page_texts).
    page_offsets = [0]
    for pt in page_texts:
        page_offsets.append(page_offsets[-1] + len(pt) + 1)  # +1 for the join newline

    headings: list[dict] = []
    page_num = getattr(reader, "get_destination_page_number", None)

    def walk(items: list, level: int = 1) -> None:
        for item in items:
            if isinstance(item, list):
                walk(item, level + 1)
                continue
            title = str(getattr(item, "title", "") or "").strip()
            if not title:
                continue
            # PDFs render headings with non-breaking spaces ("8.\xa0OSC ...");
            # normalize so TOC text, prompts, and the model's echoed paths
            # all use plain spaces (exact-match routing depends on it).
            title = title.replace("\xa0", " ")
            idx = 0
            try:
                if page_num is not None:
                    idx = int(page_num(item))  # 0-based page number
                elif getattr(item, "page", None) in reader.pages:
                    idx = reader.pages.index(item.page)
            except Exception:
                idx = 0
            headings.append(
                {
                    "title": title,
                    "level": level,
                    "offset": page_offsets[idx] if idx < len(page_offsets) else 0,
                }
            )

    walk(outline)
    return headings


class RAGSystem:
    """Main entry point for rag-kit.

    Usage:
        rag = RAGSystem()
        fid = rag.load_url("https://example.com/doc.txt")
        answer = rag.query(fid, "What is this about?")

    Auto-cleanup: When max_files is exceeded, the least recently accessed
    files are automatically deleted. Set max_files=0 for unlimited.
    """

    def __init__(
        self,
        db_path: str | None = None,
        llm_config: LLMConfig | None = None,
        default_chunk_size: int | None = None,
        default_overlap: int | None = None,
        search_threshold: float | None = None,
        max_files: int = 50,
        enable_vectors: bool = True,
        embed_backend: str = "local",
        use_cache: bool = True,
        cache_fuzzy: float | None = 0.90,
        toc_ai_headings: bool = False,
    ):
        self._storage = Storage(db_path)
        self._chunk_size = default_chunk_size or DEFAULT_CHUNK_SIZE
        self._overlap = default_overlap or DEFAULT_CHUNK_OVERLAP
        self._threshold = search_threshold or DEFAULT_SEARCH_THRESHOLD
        self._max_files = max_files
        self._use_cache = use_cache
        self._cache_fuzzy = cache_fuzzy
        # Vector index is keyed by the DATABASE (file ids are per-DB) — a
        # global index dir would collide across databases sharing file ids.
        index_dir = None
        if db_path is not None:
            index_dir = os.path.join(
                os.path.dirname(os.path.abspath(db_path)),
                os.path.basename(db_path) + ".vectors",
            )
        self._vector_index = (
            VectorIndex(embed_backend=embed_backend, index_dir=index_dir)
            if enable_vectors
            else None
        )
        # Restore the persisted index (best-effort) so a restart keeps
        # semantic search over previously loaded files instead of starting
        # empty (lexical-only) until new files are loaded.
        if self._vector_index is not None:
            try:
                self._vector_index.load("default")
            except Exception:
                pass
        self._pipeline = Pipeline(
            self._storage,
            llm_config,
            search_threshold=self._threshold,
            vector_index=self._vector_index,
            toc_ai_headings=toc_ai_headings,
        )

    def set_llm_config(self, llm_config: LLMConfig | None) -> None:
        """Set or clear the LLM configuration after construction.

        Enables upload-first, query-later pattern:
            rag = RAGSystem()
            fid = rag.load_file("doc.pdf")
            rag.set_llm_config(LLMConfig(model="gpt-4o"))
            rag.query(fid, "Summarize this.")
        """
        self._pipeline.set_llm_config(llm_config)

    def _cleanup_if_needed(self):
        """Delete least recently accessed files if over max_files limit."""
        if self._max_files <= 0:
            return
        files = self._storage.list_files(order_by="last_accessed", descending=False)
        while len(files) > self._max_files:
            f = files.pop(0)  # Oldest first
            self._storage.delete_file(f["file_id"])

    # ── Load ──────────────────────────────────────────────────────────

    def load_url(
        self,
        url: str,
        namespace: str = "default",
        chunk_size: int | None = None,
        overlap: int | None = None,
    ) -> int:
        """Fetch a URL, chunk the content, and store it.

        Args:
            url: URL to fetch.
            namespace: Logical grouping (e.g. "project-a").
            chunk_size: Max chars per chunk.
            overlap: Overlap between chunks.

        Returns:
            The new file_id.
        """
        try:
            import httpx
        except ImportError:
            raise ImportError("Install rag-kit[web] or rag-kit[llm] for URL support")

        resp = httpx.get(url, timeout=30, follow_redirects=True)
        resp.raise_for_status()
        text = _clean_text(resp.text)

        # Content hash for dedup
        content_hash = compute_content_hash(text)
        existing = self._storage.find_by_hash(content_hash, namespace)
        if existing is not None:
            return existing

        chunks = process_chunks(
            text,
            chunk_size or self._chunk_size,
            overlap or self._overlap,
        )

        # Auto-extract section mappings and TOC from raw text
        headings = _extract_headings_from_text(text)
        if headings:
            mappings = _build_section_mappings(
                headings,
                chunks,
                chunk_size or self._chunk_size,
                overlap or self._overlap,
            )
            toc_text = format_toc(mappings)
        else:
            mappings = []
            toc_text = ""

        filename = os.path.basename(url.split("?")[0]) or "unnamed"
        file_id = self._storage.create_file(
            url=url,
            file_path=None,
            filename=filename,
            source_type="url",
            content_hash=content_hash,
            chunk_size=chunk_size or self._chunk_size,
            overlap=overlap or self._overlap,
            total_chunks=len(chunks),
            chunks=chunks,
            namespace=namespace,
        )
        # Store auto-extracted section mappings and TOC
        if mappings:
            self._storage.set_section_mappings(file_id, mappings)
        if toc_text:
            self._storage.set_toc(file_id, toc_text)
        self._cleanup_if_needed()

        # Vector index the chunks (URL load)
        if self._vector_index and chunks:
            chunk_texts = [c["text"] for c in chunks]
            added = self._vector_index.add_file(file_id, chunk_texts)
            if added:
                self._vector_index.save(namespace)

        return file_id

    def load_file(
        self,
        path: str,
        namespace: str = "default",
        chunk_size: int | None = None,
        overlap: int | None = None,
    ) -> int:
        """Load a local file, chunk it, and store it.

        Args:
            path: Path to local file.
            namespace: Logical grouping (e.g. "project-a").
            chunk_size: Max chars per chunk.
            overlap: Overlap between chunks.

        Returns:
            The new file_id.
        """
        ext = os.path.splitext(path)[1].lower()
        text_exts = {
            ".txt",
            ".md",
            ".csv",
            ".json",
            ".xml",
            ".html",
            ".py",
            ".js",
            ".rs",
            ".yaml",
            ".yml",
            ".log",
        }
        reader = None
        page_texts: list[str] = []
        docling_text: str | None = None
        source_type: str | None = None
        docling_headings: list[dict] | None = None

        # Docling (optional heavy extra): when installed it handles the
        # rich formats first — real heading hierarchy, tables as
        # markdown, and bundled OCR for scans/images, so no separate OCR
        # path is needed. Formats that also have a legacy extractor fall
        # back to it if docling is missing or fails; docling-only formats
        # (xlsx, images, legacy Office, latex, email, ...) require it.
        if ext in _docling.DOCLING_EXTS and _docling.is_available():
            try:
                extracted = _docling.extract_document(path)
            except Exception as exc:
                if ext in _docling.DOCLING_ONLY_EXTS:
                    raise ValueError(
                        f"docling failed to convert {path}: {exc}. {_docling.install_hint_for(ext)}"
                    ) from exc
            else:
                cand_text = _clean_text(extracted["text"])
                if cand_text.strip() and (ext != ".pdf" or _has_meaningful_text(cand_text)):
                    docling_text = cand_text
                    source_type = _docling.source_type_for(ext)
                    docling_headings = extracted["headings"]
                # else: empty/scanned PDF docling couldn't read — fall
                # through to the legacy pypdf/OCR chain.

        if docling_text is not None:
            text = docling_text
        elif ext in text_exts:
            with open(path, encoding="utf-8", errors="ignore") as f:
                text = f.read()
            source_type = "text"
        elif ext == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError:
                raise ImportError("Install rag-kit[pdf] for PDF support")
            try:
                reader = PdfReader(path)
            except Exception as e:
                if "cryptography" in str(e) or "AES" in str(e) or "RC4" in str(e):
                    raise ImportError(
                        "This PDF is encrypted (common for datasheets/vendor docs). "
                        'Install PDF decryption support: pip install "rag-kit[pdf]" '
                        "(adds cryptography)."
                    ) from e
                raise
            page_texts = [page.extract_text() or "" for page in reader.pages]
            text = "\n".join(page_texts)
            source_type = "pdf"

            # Check if the text is meaningful (scanned PDFs extract nothing)
            if not _has_meaningful_text(text):
                # Try OCR if available
                try:
                    from pdf2image import convert_from_path
                    from tesserocr import PyTessBaseAPI

                    tessdata_path = os.path.expanduser("~/.local/share/tessdata")
                    images = convert_from_path(path)
                    ocr_lines = []
                    with PyTessBaseAPI(path=tessdata_path) as api:
                        for img in images:
                            api.SetImage(img)
                            ocr_lines.append(api.GetUTF8Text())
                    ocr_text = "\n".join(ocr_lines)
                    if _has_meaningful_text(ocr_text):
                        text = ocr_text
                    else:
                        raise ValueError(
                            "PDF appears to be scanned images with no extractable text. "
                            "Install tesseract-ocr + pytesseract + pdf2image for OCR support, "
                            "or convert the document to text manually."
                        )
                except ImportError:
                    raise ValueError(
                        f"PDF at {path} contains no extractable text (scanned document). "
                        f"Install rag-kit[ocr] for OCR support: pip install 'rag-kit[ocr]'"
                    )
        elif ext == ".docx":
            try:
                from docx import Document
            except ImportError:
                raise ImportError("Install rag-kit[docx] for DOCX support")
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs)
            source_type = "docx"
        elif ext == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                raise ImportError("Install rag-kit[pptx] for PPTX support")
            prs = Presentation(path)
            text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            source_type = "pptx"
        elif ext == ".epub":
            try:
                import ebooklib
                from ebooklib import epub
            except ImportError:
                raise ImportError("Install rag-kit[epub] for EPUB support")
            book = epub.read_epub(path)
            text = "\n".join(
                item.get_content().decode("utf-8", errors="ignore")
                for item in book.get_items()
                if item.get_type() == ebooklib.ITEM_DOCUMENT
            )
            # Strip HTML tags
            import re

            text = re.sub(r"<[^>]+>", " ", text)
            text = re.sub(r"\s+", " ", text).strip()
            source_type = "epub"
        elif ext == ".odt":
            try:
                from odf import teletype, text
                from odf.opendocument import load
            except ImportError:
                raise ImportError("Install rag-kit[odt] for ODT support")
            doc = load(path)
            paras = doc.getElementsByType(text.P)
            text = "\n".join(teletype.retrieveText(p) for p in paras)
            source_type = "odt"
        elif ext == ".rtf":
            try:
                from striprtf.striprtf import rtf_to_text
            except ImportError:
                raise ImportError("Install rag-kit[rtf] for RTF support")
            with open(path, encoding="utf-8", errors="ignore") as f:
                text = rtf_to_text(f.read())
            source_type = "rtf"
        else:
            if ext in _docling.DOCLING_ONLY_EXTS:
                raise ImportError(
                    f"Format {ext} requires docling. {_docling.install_hint_for(ext)}"
                )
            raise ValueError(f"Unsupported file type: {ext}")

        # Clean surrogates from extracted text before hashing or chunking
        text = _clean_text(text)

        # Content hash for dedup
        content_hash = compute_content_hash(text)
        existing = self._storage.find_by_hash(content_hash, namespace)
        if existing is not None:
            return existing

        chunks = process_chunks(
            text,
            chunk_size or self._chunk_size,
            overlap or self._overlap,
        )

        # Auto-extract section mappings and TOC. Docling's real heading
        # hierarchy wins when present; PDFs with an embedded outline
        # (datasheets, manuals) get a structured TOC from the bookmarks;
        # everything else uses the text heuristic.
        headings: list[dict] = []
        if docling_headings:
            headings = docling_headings
        elif source_type == "pdf" and reader is not None:
            headings = _pdf_outline_headings(reader, page_texts)
        if not headings:
            headings = _extract_headings_from_text(text)
        if headings:
            mappings = _build_section_mappings(
                headings,
                chunks,
                chunk_size or self._chunk_size,
                overlap or self._overlap,
            )
            toc_text = format_toc(mappings)
        else:
            mappings = []
            toc_text = ""

        filename = os.path.basename(path)
        file_id = self._storage.create_file(
            url=None,
            file_path=path,
            filename=filename,
            source_type=source_type,
            content_hash=content_hash,
            chunk_size=chunk_size or self._chunk_size,
            overlap=overlap or self._overlap,
            total_chunks=len(chunks),
            chunks=chunks,
            namespace=namespace,
        )
        # Store auto-extracted section mappings and TOC
        if mappings:
            self._storage.set_section_mappings(file_id, mappings)
        if toc_text:
            self._storage.set_toc(file_id, toc_text)
        self._cleanup_if_needed()

        # Vector index the chunks (file load)
        if self._vector_index and chunks:
            chunk_texts = [c["text"] for c in chunks]
            added = self._vector_index.add_file(file_id, chunk_texts)
            if added:
                self._vector_index.save(namespace)

        return file_id

    # ── Query ─────────────────────────────────────────────────────────

    def query(
        self,
        file_id_or_question: int | str,
        question: str | None = None,
        namespace: str | None = None,
        llm_config: LLMConfig | None = None,
        toc_first: bool = False,
        terse: bool = False,
        expand_terms: bool = True,
        conversation: str | None = None,
    ) -> QueryResult:
        """Ask a question about a loaded document.

        Two calling modes:
        1. By file_id: rag.query(file_id, "question")
        2. By namespace: rag.query("question", namespace="project-a")

        Args:
            file_id_or_question: File ID or question string.
            question: Question (if first arg is file_id).
            namespace: Namespace to search (if querying by namespace).
            llm_config: Optional per-query LLM config override.
            toc_first: Use TOC-first retrieval (relevant for file_id queries).
            expand_terms: In TOC-first mode, spawn 3-7 search terms after
                reading the TOC and search them in parallel (default True).
                Measured: recovers hard questions at ~3.8x prompt cost — keep
                on for identifier-heavy corpora, off for cheap factoid use.

        Returns:
            QueryResult with .answer (str) and .citations (list[dict]).
        """
        cache_ctx = None
        # Conversational turns (conversation passed) skip the cache: the
        # answer depends on the thread, and the key is question-only.
        if self._use_cache and conversation is None:
            if question is not None:
                scope = f"file:{file_id_or_question}"
                qtext = question
            elif namespace is not None:
                scope = f"ns:{namespace}"
                qtext = str(file_id_or_question)
            else:
                scope = "ns:__all__"
                qtext = str(file_id_or_question)
            qnorm = _qnorm(qtext)
            hit = self._storage.cache_lookup(scope, qnorm, self._cache_fuzzy)
            if hit is not None:
                metrics = {"cached": True, "cache_hits": hit["hits"]}
                if "fuzzy_ratio" in hit:
                    metrics["cache_fuzzy_ratio"] = hit["fuzzy_ratio"]
                return QueryResult(
                    answer=hit["answer"], citations=hit["citations"], metrics=metrics
                )
            cache_ctx = (scope, qnorm, qtext)

        if question is not None:
            # Mode 1: file_id + question
            if toc_first:
                answer, citations = self._pipeline.query_toc_first(
                    file_id=file_id_or_question,
                    question=question,
                    llm_config=llm_config,
                    expand_terms=expand_terms,
                    conversation=conversation,
                )
            else:
                answer, citations = self._pipeline.query(
                    file_id=file_id_or_question,
                    question=question,
                    llm_config=llm_config,
                    terse=terse,
                    conversation=conversation,
                )
        elif namespace is not None:
            # Mode 2: question + namespace (cross-file search)
            answer, citations = self._pipeline.query_by_namespace(
                question=str(file_id_or_question),
                namespace=namespace,
                llm_config=llm_config,
            )
        else:
            # Mode 3: question only (cross-file, all namespaces)
            answer, citations = self._pipeline.query_by_namespace(
                question=str(file_id_or_question),
                namespace=None,
                llm_config=llm_config,
            )

        if cache_ctx is not None and answer.strip() and citations:
            self._storage.cache_put(cache_ctx[0], cache_ctx[1], cache_ctx[2], answer, citations)
        return QueryResult(answer=answer, citations=citations)

    async def aquery(
        self,
        file_id_or_question: int | str,
        question: str | None = None,
        namespace: str | None = None,
        llm_config: LLMConfig | None = None,
        toc_first: bool = False,
        terse: bool = False,
        expand_terms: bool = True,
        conversation: str | None = None,
    ) -> QueryResult:
        """Async query — same contract as query() with async LLM synthesis.

        File-ID mode uses the async pipeline (awaitable, non-blocking).
        Namespace/question-only modes fall back to the sync pipeline.
        """
        cache_ctx = None
        # Conversational turns skip the cache (answer depends on thread).
        if self._use_cache and conversation is None:
            if question is not None:
                scope = f"file:{file_id_or_question}"
                qtext = question
            elif namespace is not None:
                scope = f"ns:{namespace}"
                qtext = str(file_id_or_question)
            else:
                scope = "ns:__all__"
                qtext = str(file_id_or_question)
            qnorm = _qnorm(qtext)
            hit = self._storage.cache_lookup(scope, qnorm, self._cache_fuzzy)
            if hit is not None:
                metrics = {"cached": True, "cache_hits": hit["hits"]}
                if "fuzzy_ratio" in hit:
                    metrics["cache_fuzzy_ratio"] = hit["fuzzy_ratio"]
                return QueryResult(
                    answer=hit["answer"], citations=hit["citations"], metrics=metrics
                )
            cache_ctx = (scope, qnorm, qtext)

        if question is not None:
            if toc_first:
                answer, citations = self._pipeline.query_toc_first(
                    file_id=file_id_or_question,
                    question=question,
                    llm_config=llm_config,
                    expand_terms=expand_terms,
                    conversation=conversation,
                )
            else:
                answer, citations = await self._pipeline.aquery(
                    file_id=file_id_or_question,
                    question=question,
                    llm_config=llm_config,
                    terse=terse,
                    conversation=conversation,
                )
        elif namespace is not None:
            answer, citations = self._pipeline.query_by_namespace(
                question=str(file_id_or_question),
                namespace=namespace,
                llm_config=llm_config,
            )
        else:
            answer, citations = self._pipeline.query_by_namespace(
                question=str(file_id_or_question),
                namespace=None,
                llm_config=llm_config,
            )

        if cache_ctx is not None and answer.strip() and citations:
            self._storage.cache_put(cache_ctx[0], cache_ctx[1], cache_ctx[2], answer, citations)
        return QueryResult(answer=answer, citations=citations)

    def query_agentic(
        self,
        file_id: int,
        question: str,
        llm_config: LLMConfig | None = None,
        max_turns: int = 10,
        searcher_model: str | None = None,
        planner_model: str | None = None,
    ) -> QueryResult:
        """Agentic RAG: LLM searches the document itself using a search tool.

        The LLM decides what to search for, iterates with a search_document
        tool, and synthesises the answer from the results it gathers.

        Args:
            file_id: ID of the loaded file.
            question: Question to ask.
            llm_config: Optional per-query LLM config override.
            max_turns: Max tool-calling iterations (default 10).
            searcher_model: Model for the executor stage.
            planner_model: Model for the planner stage.

        Returns:
            QueryResult with .answer and .citations.
        """
        answer, citations, metrics = self._pipeline.query_agentic(
            file_id=file_id,
            question=question,
            llm_config=llm_config,
            max_turns=max_turns,
            searcher_model=searcher_model,
            planner_model=planner_model,
        )
        return QueryResult(answer=answer, citations=citations, metrics=metrics)

    def query_loop(
        self,
        file_id: int,
        question: str,
        llm_config: LLMConfig | None = None,
        max_loops: int = 4,
        verifier_model: str | None = None,
        conversation: str | None = None,
    ) -> QueryResult:
        """Iterative retrieval loop with a cheap sufficiency verifier.

        Deterministic Self-RAG-style loop: search the original question,
        ask a cheap verifier whether the collected excerpts are enough to
        answer, and if not search the terms it suggests — repeat until
        sufficient, max_loops, or no new evidence. Cheaper and more
        predictable than query_agentic (no planner, no tool-calling).

        Args:
            file_id: ID of the loaded file.
            question: Question to ask.
            llm_config: Optional per-query LLM config override (synthesis).
            max_loops: Max retrieval rounds after the initial search
                (default 4; hard-capped at 10 — the loop always concludes
                with a final answer once the cap is reached).
            verifier_model: Model id for the sufficiency verifier
                (defaults to the router model).

        Returns:
            QueryResult with .answer, .citations and .metrics (stop_reason,
            loops, verifier_calls, chunks_found).
        """
        answer, citations, metrics = self._pipeline.query_loop(
            file_id=file_id,
            question=question,
            llm_config=llm_config,
            max_loops=max_loops,
            verifier_model=verifier_model,
            conversation=conversation,
        )
        return QueryResult(answer=answer, citations=citations, metrics=metrics)

    # ── Search ────────────────────────────────────────────────────────

    def search(
        self,
        query: str,
        file_id: int | None = None,
        namespace: str | None = None,
    ) -> list[dict[str, Any]]:
        """Direct keyword search without LLM.

        Uses hybrid (vector + FTS5) when vector index is available,
        falls back to fuzzy + FTS5.

        Returns matching chunks sorted by relevance.
        """
        return search(
            storage=self._storage,
            query=query,
            file_id=file_id,
            namespace=namespace,
            top_k=20,
            threshold=self._threshold,
            vector_index=self._vector_index,
        )

    def algorithmic_search(
        self, file_id: int, question: str, top_k: int = 8
    ) -> list[dict[str, Any]]:
        """Algorithmic retrieval WITHOUT synthesis — the backend for the
        chat's search_documents tool.

        Runs the TOC-first engine: TOC heading selection → term expansion
        → parallel section-scoped + hybrid + lexical search → fusion →
        cross-encoder rerank (see Pipeline._algorithmic_retrieve). The
        chat model writes the answer; this returns the best chunks (with
        section names) so the model never sees raw retrieval internals.
        """
        chunks = self._pipeline._algorithmic_retrieve(file_id, question)
        return chunks[:top_k]

    def loop_retrieve(
        self, file_id: int, question: str, max_loops: int = 2, top_k: int = 8
    ) -> list[dict[str, Any]]:
        """Verifier-driven iterative retrieval WITHOUT synthesis — the
        'loop' algorithm as a chat tool backend. Search → cheap verifier
        asks for more terms if evidence is insufficient → repeat. Returns
        the collected chunks (with section names)."""
        collected, _metrics = self._pipeline.retrieve_loop(
            question,
            file_id=file_id,
            max_loops=max_loops,
        )
        mappings = self._storage.get_section_mappings(file_id) or []
        if mappings:
            try:
                collected = self._pipeline._expand_context(collected, mappings, file_id)
            except Exception:
                pass
        return collected[:top_k]

    # ── Chunk / TOC access ────────────────────────────────────────────

    def get_chunk(self, file_id: int, index: int) -> dict | None:
        return self._storage.get_chunk(file_id, index)

    def get_toc(self, file_id: int) -> str | None:
        return self._storage.get_toc(file_id)

    def update_toc(self, file_id: int, toc_text: str) -> bool:
        return self._storage.set_toc(file_id, toc_text)

    # ── File management ───────────────────────────────────────────────

    def list(self, namespace: str | None = None) -> list[dict]:
        """List files, optionally filtered by namespace."""
        return self._storage.list_files(namespace=namespace)

    def delete_file(self, file_id: int) -> bool:
        # Remove from vector index first
        if self._vector_index:
            info = self._storage.get_file(file_id)
            if info:
                self._vector_index.remove_file(file_id, info.get("total_chunks", 0))
                namespace = info.get("namespace", "default")
                self._vector_index.save(namespace)
        return self._storage.delete_file(file_id)

    def stats(self) -> dict:
        return self._storage.stats()
